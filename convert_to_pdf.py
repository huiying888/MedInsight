import os
from pathlib import Path
from PIL import Image
import pandas as pd
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Preformatted, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from docx import Document
from pptx import Presentation
import io
import shutil

class FileConverter:
    def __init__(self, input_folder=None, output_folder=None):
        if input_folder:
            self.input_folder = Path(input_folder)
        self.output_folder = Path(output_folder) if output_folder else "pdf_output"
        self.output_folder.mkdir(exist_ok=True)
        
    def convert_image_to_pdf(self, input_path, output_path):
        """Convert image files (PNG, JPG, JPEG) to PDF"""
        try:
            image = Image.open(input_path)
            
            # Convert to RGB if necessary
            if image.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', image.size, (255, 255, 255))
                if image.mode == 'P':
                    image = image.convert('RGBA')
                background.paste(image, mask=image.split()[-1] if image.mode in ('RGBA', 'LA') else None)
                image = background
            
            image.save(output_path, 'PDF', resolution=100.0)
            print(f"✓ Converted: {input_path.name} -> {output_path.name}")
            return True
        except Exception as e:
            print(f"✗ Error converting {input_path.name}: {str(e)}")
            return False
    
    def convert_docx_to_pdf(self, input_path, output_path):
        """Convert DOCX to PDF with text, tables, and images"""
        try:
            doc = Document(input_path)
            pdf = SimpleDocTemplate(str(output_path), pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            # Process paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    style = styles['Heading1'] if para.style.name.startswith('Heading') else styles['BodyText']
                    elements.append(Paragraph(para.text, style))
                    elements.append(Spacer(1, 6))
            
            # Process tables
            for table in doc.tables:
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                if table_data:
                    t = Table(table_data)
                    t.setStyle(TableStyle([
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ]))
                    elements.append(t)
                    elements.append(Spacer(1, 12))
            
            # Extract images from document relationships
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image = rel.target_part
                        image_bytes = image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        
                        # Save to buffer
                        img_buffer = io.BytesIO()
                        img.save(img_buffer, format="PNG")
                        img_buffer.seek(0)
                        
                        # Original size (pixels)
                        img_width_px, img_height_px = img.size
                        
                        # Convert to points
                        dpi = img.info.get("dpi", (96, 96))[0]  # fallback to 96 dpi
                        img_width_pt = img_width_px * 72 / dpi
                        img_height_pt = img_height_px * 72 / dpi
                        
                        # Fit within page content area
                        max_width = A4[0] - 3*inch   # leave 1.5 inch margin left/right
                        max_height = A4[1] - 3*inch  # leave 1.5 inch margin top/bottom
                        scale = min(1, max_width / img_width_pt, max_height / img_height_pt)
                        
                        rl_img = RLImage(
                            img_buffer,
                            width=img_width_pt * scale,
                            height=img_height_pt * scale
                        )
                        
                        elements.append(rl_img)
                        elements.append(Spacer(1, 12))
                    except Exception as img_error:
                        print(f"  Warning: Could not extract image - {img_error}")
            
            pdf.build(elements)
            print(f"✓ Converted: {input_path.name} -> {output_path.name}")
            return True
        except Exception as e:
            print(f"✗ Error converting {input_path.name}: {str(e)}")
            return False
    
    def convert_xlsx_to_pdf(self, input_path, output_path):
        """Convert Excel to PDF"""
        try:
            df = pd.read_excel(input_path, sheet_name=None)
            doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            for sheet_name, data in df.items():
                elements.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                elements.append(Spacer(1, 12))
                
                table_data = [data.columns.tolist()] + data.values.tolist()
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                elements.append(table)
                elements.append(Spacer(1, 24))
            
            doc.build(elements)
            print(f"✓ Converted: {input_path.name} -> {output_path.name}")
            return True
        except Exception as e:
            print(f"✗ Error converting {input_path.name}: {str(e)}")
            return False
    
    def convert_ppt_to_pdf(self, input_path, output_path):
        """Convert PPT/PPTX to PDF with text, tables, and images"""
        try:
            prs = Presentation(input_path)
            pdf = SimpleDocTemplate(str(output_path), pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide number
                elements.append(Paragraph(f"Slide {slide_num}", styles['Heading1']))
                elements.append(Spacer(1, 6))
                
                # Extract text and images from shapes
                for shape in slide.shapes:
                    # Extract text
                    if hasattr(shape, "text") and shape.text.strip():
                        style = styles['Heading2'] if hasattr(shape, "is_placeholder") and shape.is_placeholder else styles['BodyText']
                        elements.append(Paragraph(shape.text, style))
                        elements.append(Spacer(1, 6))
                    
                    # Extract images
                    if shape.shape_type == 13:  # picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            
                            # Save to buffer
                            img_buffer = io.BytesIO()
                            img.save(img_buffer, format="PNG")
                            img_buffer.seek(0)
                            
                            # Original size (pixels)
                            img_width_px, img_height_px = img.size
                            
                            # Convert to points
                            dpi = img.info.get("dpi", (96, 96))[0]  # fallback to 96 dpi
                            img_width_pt = img_width_px * 72 / dpi
                            img_height_pt = img_height_px * 72 / dpi
                            
                            # Fit within page content area
                            max_width = A4[0] - 3*inch   # leave 1.5 inch margin left/right
                            max_height = A4[1] - 3*inch  # leave 1.5 inch margin top/bottom
                            scale = min(1, max_width / img_width_pt, max_height / img_height_pt)
                            
                            rl_img = RLImage(
                                img_buffer,
                                width=img_width_pt * scale,
                                height=img_height_pt * scale
                            )
                            
                            elements.append(rl_img)
                            elements.append(Spacer(1, 12))
                        except Exception as img_error:
                            print(f"  Warning: Could not extract image from slide {slide_num} - {img_error}")
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                        if table_data:
                            t = Table(table_data)
                            t.setStyle(TableStyle([
                                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ]))
                            elements.append(t)
                            elements.append(Spacer(1, 12))
                
                # Page break between slides
                if slide_num < len(prs.slides):
                    elements.append(PageBreak())
            
            pdf.build(elements)
            print(f"✓ Converted: {input_path.name} -> {output_path.name}")
            return True
        except Exception as e:
            print(f"✗ Error converting {input_path.name}: {str(e)}")
            return False
    
    def convert_csv_to_pdf(self, input_path, output_path):
        """Convert CSV to PDF"""
        try:
            df = pd.read_csv(input_path)
            doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            elements = []
            
            table_data = [df.columns.tolist()] + df.values.tolist()
            table = Table(table_data, repeatRows=1)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            elements.append(table)
            
            doc.build(elements)
            print(f"✓ Converted: {input_path.name} -> {output_path.name}")
            return True
        except Exception as e:
            print(f"✗ Error converting {input_path.name}: {str(e)}")
            return False
        
    def convert_txt_to_pdf(self, input_path, output_path):
        """Convert TXT file to PDF"""
        try:
            # Read text file
            with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()

            # Use Preformatted so spacing, tabs, and newlines are preserved
            pre_style = styles["Normal"]
            pre_style.fontName = "Courier"   # monospaced font for TXT
            pre_style.fontSize = 10

            elements.append(Preformatted(content, pre_style))
            
            doc.build(elements)
            print(f"✓ Converted: {Path(input_path).name} -> {Path(output_path).name}")
            return True

        except Exception as e:
            print(f"✗ Error converting {Path(input_path).name}: {str(e)}")
            return False
        
    def convert_file(self, file_path, output_path=None):
        """Convert a single file to PDF based on its extension"""
        file_path = Path(file_path)
        if not output_path:
            output_path = self.output_folder / f"{file_path.stem}.pdf"
        
        ext = file_path.suffix.lower()
        
        if ext == ".pdf":
            shutil.copy(file_path, output_path)
            print(f"✓ Copied PDF: {file_path.name} -> {output_path.name}")
            return True
        elif ext in ['.png', '.jpg', '.jpeg']:
            return self.convert_image_to_pdf(file_path, output_path)
        elif ext == '.docx':
            return self.convert_docx_to_pdf(file_path, output_path)
        elif ext == '.xlsx':
            return self.convert_xlsx_to_pdf(file_path, output_path)
        elif ext in ['.ppt', '.pptx']:
            return self.convert_ppt_to_pdf(file_path, output_path)
        elif ext == '.csv':
            return self.convert_csv_to_pdf(file_path, output_path)
        elif ext == '.txt':
            return self.convert_txt_to_pdf(file_path, output_path)
        else:
            print(f"⊘ Unsupported format: {file_path.name}")
            return False
    
    def convert_all(self):
        """Convert all supported files in the input folder"""
        supported_extensions = ['.pdf', '.png', '.jpg', '.jpeg', '.docx', '.xlsx', '.ppt', '.pptx', '.csv', '.txt']
        files = [f for f in self.input_folder.iterdir() if f.suffix.lower() in supported_extensions]
        
        if not files:
            print("No supported files found in the input folder.")
            return
        
        print(f"\nFound {len(files)} file(s) to convert.\n")
        
        successful = 0
        failed = 0
        
        for file in files:
            if self.convert_file(file):
                successful += 1
            else:
                failed += 1
        
        print(f"\n{'='*50}")
        print(f"Conversion complete!")
        print(f"Successful: {successful}")
        print(f"Failed: {failed}")
        print(f"Output folder: {self.output_folder}")
        print(f"{'='*50}")


# Example usage
if __name__ == "__main__":
    # Set your input folder path
    input_folder = "input_folder"
    
    # Optional: specify output folder (default is 'pdf_output' inside input folder)
    output_folder = "output_folder"
    
    # Create converter and convert all files
    converter = FileConverter(input_folder, output_folder)
    converter.convert_all()
    
    # Or convert a single file
    # converter.convert_file("path/to/single/file.png")